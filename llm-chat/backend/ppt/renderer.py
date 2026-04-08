"""
PPT 渲染器 — 将结构化 JSON 渲染为 .pptx 文件

面向对象设计：
  PptRenderer        — 渲染器主类，管理 Presentation 对象
  SlideBuilder       — 基类，定义幻灯片构建接口
  TitleSlideBuilder  — 封面页
  ContentSlideBuilder — 要点列表页
  TwoColumnBuilder   — 双栏对比页
  ImageTextBuilder   — 图文混排页
  ChartSlideBuilder  — 图表页
  QuoteSlideBuilder  — 引用页
  EndingSlideBuilder — 结束页
"""
import io
import logging
from abc import ABC, abstractmethod
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE

from ppt.themes import Theme, get_theme

logger = logging.getLogger("ppt.renderer")


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """#RRGGBB → RGBColor"""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ═══════════════════════════════════════════════════════════════════════════════
# SlideBuilder — 幻灯片构建基类
# ═══════════════════════════════════════════════════════════════════════════════

class SlideBuilder(ABC):
    """幻灯片构建器基类。"""

    def __init__(self, prs: Presentation, theme: Theme):
        self.prs = prs
        self.theme = theme
        self.slide_width = prs.slide_width
        self.slide_height = prs.slide_height

    @abstractmethod
    def build(self, slide_data: dict) -> None:
        """根据 slide_data 构建一张幻灯片。"""

    def _add_blank_slide(self):
        layout = self.prs.slide_layouts[6]  # blank layout
        return self.prs.slides.add_slide(layout)

    def _add_textbox(self, slide, left, top, width, height, text: str,
                     font_size: int = 18, font_color: str = "", bold: bool = False,
                     alignment=PP_ALIGN.LEFT, font_name: str = ""):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = alignment
        run = p.runs[0] if p.runs else p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.name = font_name or self.theme.body_font
        if font_color:
            run.font.color.rgb = _hex_to_rgb(font_color)
        return txBox

    def _add_bullets(self, slide, left, top, width, height, bullets: list[str],
                     font_size: int = 16, font_color: str = ""):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        color = font_color or self.theme.colors.text
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"  {bullet}"
            p.space_before = Pt(6)
            p.space_after = Pt(4)
            p.level = 0
            run = p.runs[0] if p.runs else p.add_run()
            run.text = f"  {bullet}"
            run.font.size = Pt(font_size)
            run.font.name = self.theme.body_font
            run.font.color.rgb = _hex_to_rgb(color)
        return txBox

    def _add_accent_line(self, slide, left, top, width):
        from pptx.util import Pt as PtUtil
        shape = slide.shapes.add_shape(
            1,  # rectangle
            left, top, width, Pt(3),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(self.theme.colors.accent)
        shape.line.fill.background()
        return shape

    def _set_slide_bg(self, slide, color: str = ""):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(color or self.theme.colors.bg)


# ═══════════════════════════════════════════════════════════════════════════════
# 具体 Builder 实现
# ═══════════════════════════════════════════════════════════════════════════════

class TitleSlideBuilder(SlideBuilder):
    """封面页：大标题 + 副标题 + 装饰线。"""

    def build(self, data: dict) -> None:
        slide = self._add_blank_slide()
        self._set_slide_bg(slide)

        # 装饰线
        self._add_accent_line(slide, Inches(1.5), Inches(3.2), Inches(3))

        # 标题
        self._add_textbox(
            slide, Inches(1.5), Inches(1.8), Inches(7), Inches(1.5),
            data.get("title", ""), font_size=36, bold=True,
            font_color=self.theme.colors.primary,
            font_name=self.theme.title_font,
            alignment=PP_ALIGN.LEFT,
        )
        # 副标题
        subtitle = data.get("subtitle", "")
        if subtitle:
            self._add_textbox(
                slide, Inches(1.5), Inches(3.5), Inches(7), Inches(0.8),
                subtitle, font_size=18,
                font_color=self.theme.colors.secondary,
            )


class ContentSlideBuilder(SlideBuilder):
    """标题 + 要点列表页。"""

    def build(self, data: dict) -> None:
        slide = self._add_blank_slide()
        self._set_slide_bg(slide)

        # 标题
        self._add_textbox(
            slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.8),
            data.get("title", ""), font_size=28, bold=True,
            font_color=self.theme.colors.primary,
            font_name=self.theme.title_font,
        )
        # 装饰线
        self._add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5))

        # 要点列表
        bullets = data.get("bullets", [])
        if bullets:
            self._add_bullets(
                slide, Inches(0.8), Inches(1.5), Inches(8.4), Inches(5),
                bullets, font_size=18,
            )


class TwoColumnBuilder(SlideBuilder):
    """双栏对比页。"""

    def build(self, data: dict) -> None:
        slide = self._add_blank_slide()
        self._set_slide_bg(slide)

        # 标题
        self._add_textbox(
            slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.8),
            data.get("title", ""), font_size=28, bold=True,
            font_color=self.theme.colors.primary,
            font_name=self.theme.title_font,
        )
        self._add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(1.5))

        left_data = data.get("left", {})
        right_data = data.get("right", {})

        # 左栏标题
        if left_data.get("heading"):
            self._add_textbox(
                slide, Inches(0.8), Inches(1.5), Inches(4), Inches(0.5),
                left_data["heading"], font_size=20, bold=True,
                font_color=self.theme.colors.accent,
            )
        if left_data.get("bullets"):
            self._add_bullets(
                slide, Inches(0.8), Inches(2.1), Inches(4), Inches(4.5),
                left_data["bullets"], font_size=16,
            )

        # 右栏标题
        if right_data.get("heading"):
            self._add_textbox(
                slide, Inches(5.2), Inches(1.5), Inches(4), Inches(0.5),
                right_data["heading"], font_size=20, bold=True,
                font_color=self.theme.colors.accent,
            )
        if right_data.get("bullets"):
            self._add_bullets(
                slide, Inches(5.2), Inches(2.1), Inches(4), Inches(4.5),
                right_data["bullets"], font_size=16,
            )


class ImageTextBuilder(SlideBuilder):
    """图文混排页（图片占左半，文字占右半）。"""

    def build(self, data: dict) -> None:
        slide = self._add_blank_slide()
        self._set_slide_bg(slide)

        # 标题
        self._add_textbox(
            slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.8),
            data.get("title", ""), font_size=28, bold=True,
            font_color=self.theme.colors.primary,
            font_name=self.theme.title_font,
        )

        # 图片占位（灰色矩形 + 提示文字）
        image_data = data.get("image_data")  # base64 bytes or None
        if image_data and isinstance(image_data, bytes):
            img_stream = io.BytesIO(image_data)
            slide.shapes.add_picture(img_stream, Inches(0.8), Inches(1.4), Inches(4.2), Inches(4.8))
        else:
            # 占位矩形
            shape = slide.shapes.add_shape(1, Inches(0.8), Inches(1.4), Inches(4.2), Inches(4.8))
            shape.fill.solid()
            shape.fill.fore_color.rgb = _hex_to_rgb("#F0F0F0")
            shape.line.color.rgb = _hex_to_rgb("#E0E0E0")
            # 占位文字
            prompt = data.get("image_prompt", "图片")
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"[{prompt}]"
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.size = Pt(12)
            run.font.color.rgb = _hex_to_rgb("#999999")

        # 右侧文字
        text = data.get("text", "")
        if text:
            self._add_textbox(
                slide, Inches(5.3), Inches(1.4), Inches(4), Inches(4.8),
                text, font_size=16, font_color=self.theme.colors.text,
            )


class ChartSlideBuilder(SlideBuilder):
    """图表页（柱形图/饼图/折线图）。"""

    def build(self, data: dict) -> None:
        slide = self._add_blank_slide()
        self._set_slide_bg(slide)

        # 标题
        self._add_textbox(
            slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.8),
            data.get("title", ""), font_size=28, bold=True,
            font_color=self.theme.colors.primary,
            font_name=self.theme.title_font,
        )

        chart_data_raw = data.get("chart_data", {})
        chart_type_str = chart_data_raw.get("type", "bar")
        categories = chart_data_raw.get("categories", [])
        series_list = chart_data_raw.get("series", [])

        if not categories or not series_list:
            self._add_textbox(
                slide, Inches(2), Inches(3), Inches(6), Inches(1),
                "[图表数据不足]", font_size=16, font_color="#999999",
                alignment=PP_ALIGN.CENTER,
            )
            return

        chart_type_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "pie": XL_CHART_TYPE.PIE,
            "line": XL_CHART_TYPE.LINE_MARKERS,
        }
        chart_type = chart_type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

        from pptx.chart.data import CategoryChartData
        cd = CategoryChartData()
        cd.categories = categories
        for s in series_list:
            cd.add_series(s.get("name", ""), s.get("values", []))

        slide.shapes.add_chart(
            chart_type, Inches(1), Inches(1.5), Inches(8), Inches(5), cd,
        )


class QuoteSlideBuilder(SlideBuilder):
    """引用/金句页。"""

    def build(self, data: dict) -> None:
        slide = self._add_blank_slide()
        self._set_slide_bg(slide)

        quote = data.get("quote", data.get("text", ""))
        author = data.get("author", "")

        # 大号引号
        self._add_textbox(
            slide, Inches(1.5), Inches(1.5), Inches(1), Inches(1),
            "\u201C", font_size=72, font_color=self.theme.colors.accent,
            bold=True,
        )

        # 引文
        self._add_textbox(
            slide, Inches(1.5), Inches(2.5), Inches(7), Inches(2.5),
            quote, font_size=24, font_color=self.theme.colors.text,
            alignment=PP_ALIGN.LEFT,
        )

        # 作者
        if author:
            self._add_textbox(
                slide, Inches(1.5), Inches(5.2), Inches(7), Inches(0.5),
                f"—— {author}", font_size=16,
                font_color=self.theme.colors.secondary,
                alignment=PP_ALIGN.RIGHT,
            )


class EndingSlideBuilder(SlideBuilder):
    """结束页。"""

    def build(self, data: dict) -> None:
        slide = self._add_blank_slide()
        self._set_slide_bg(slide)

        title = data.get("title", "Thanks")
        subtitle = data.get("subtitle", "")

        self._add_accent_line(slide, Inches(3.5), Inches(3), Inches(3))

        self._add_textbox(
            slide, Inches(1), Inches(2), Inches(8), Inches(1),
            title, font_size=44, bold=True,
            font_color=self.theme.colors.primary,
            font_name=self.theme.title_font,
            alignment=PP_ALIGN.CENTER,
        )
        if subtitle:
            self._add_textbox(
                slide, Inches(1), Inches(3.3), Inches(8), Inches(0.6),
                subtitle, font_size=16,
                font_color=self.theme.colors.secondary,
                alignment=PP_ALIGN.CENTER,
            )


# ═══════════════════════════════════════════════════════════════════════════════
# 布局注册表
# ═══════════════════════════════════════════════════════════════════════════════

LAYOUT_BUILDERS: dict[str, type[SlideBuilder]] = {
    "title": TitleSlideBuilder,
    "content": ContentSlideBuilder,
    "two_column": TwoColumnBuilder,
    "image_text": ImageTextBuilder,
    "chart": ChartSlideBuilder,
    "quote": QuoteSlideBuilder,
    "ending": EndingSlideBuilder,
}


# ═══════════════════════════════════════════════════════════════════════════════
# PptRenderer — 渲染器主类
# ═══════════════════════════════════════════════════════════════════════════════

class PptRenderer:
    """
    将结构化 JSON 渲染为 .pptx 文件。

    用法：
        renderer = PptRenderer(theme="tech_blue")
        pptx_bytes = renderer.render(slides_json)
    """

    def __init__(self, theme: str = "tech_blue"):
        self.theme = get_theme(theme)
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

    def render(self, ppt_data: dict) -> bytes:
        """渲染完整 PPT，返回 .pptx 文件的 bytes。"""
        slides = ppt_data.get("slides", [])

        for slide_data in slides:
            layout = slide_data.get("layout", "content")
            builder_cls = LAYOUT_BUILDERS.get(layout, ContentSlideBuilder)
            builder = builder_cls(self.prs, self.theme)
            try:
                builder.build(slide_data)
            except Exception as exc:
                logger.warning("幻灯片渲染失败 | layout=%s | error=%s", layout, exc)
                # 降级：用 content 布局展示错误信息
                fallback = ContentSlideBuilder(self.prs, self.theme)
                fallback.build({
                    "title": slide_data.get("title", "渲染失败"),
                    "bullets": [f"错误: {exc}"],
                })

        buf = io.BytesIO()
        self.prs.save(buf)
        return buf.getvalue()

    def render_to_file(self, ppt_data: dict, path: str) -> str:
        """渲染并保存到文件，返回文件路径。"""
        data = self.render(ppt_data)
        with open(path, "wb") as f:
            f.write(data)
        return path
