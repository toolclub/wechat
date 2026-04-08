"""
PPT 脚本生成器 — 生成可在沙箱中独立运行的 Python 渲染脚本

HTML-first 方案：
  1. 将每页 HTML 写入临时文件
  2. 使用 Playwright 截图每页 HTML（960x720）
  3. 将截图嵌入 PPT 幻灯片（全页铺满）

兼容旧方案（generate_render_script）保留，但新流程优先使用 generate_html_render_script。
"""
import json

from ppt.themes import get_theme


def generate_html_render_script(slides_html: list[str], output_filename: str) -> str:
    """
    生成 HTML→截图→PPT 的自包含渲染脚本。

    脚本在沙箱中执行，需要 playwright + python-pptx。
    """
    # 将 HTML 列表序列化为 JSON（处理转义）
    html_json = json.dumps(slides_html, ensure_ascii=False)

    return f'''#!/usr/bin/env python3
"""HTML→截图→PPT 渲染脚本（沙箱执行）"""
import json, os, sys

SLIDES_HTML = json.loads("""{html_json}""")
OUTPUT = "{output_filename}"

print(f"  📄 共 {{len(SLIDES_HTML)}} 页幻灯片")

# ── 写入 HTML 文件 ──
html_dir = "slide_html"
os.makedirs(html_dir, exist_ok=True)
for i, html in enumerate(SLIDES_HTML):
    path = os.path.join(html_dir, f"slide_{{i}}.html")
    with open(path, "w", encoding="utf-8") as f:
        f.write(html)
    print(f"  ✅ 第 {{i+1}} 页 HTML 已写入")

# ── Playwright 截图 ──
print("  📸 正在启动浏览器截图...")
from playwright.sync_api import sync_playwright

screenshots = []
with sync_playwright() as p:
    browser = p.chromium.launch()
    page = browser.new_page(viewport={{"width": 960, "height": 720}})

    for i in range(len(SLIDES_HTML)):
        html_path = os.path.abspath(os.path.join(html_dir, f"slide_{{i}}.html"))
        page.goto(f"file://{{html_path}}")
        page.wait_for_timeout(300)  # 等待渲染
        img_path = f"slide_{{i}}.png"
        page.screenshot(path=img_path)
        screenshots.append(img_path)
        print(f"  📸 第 {{i+1}} 页截图完成")

    browser.close()

# ── 组装 PPT ──
print("  📦 正在组装 PPT...")
from pptx import Presentation
from pptx.util import Inches, Emu

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

for i, img_path in enumerate(screenshots):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    # 全页铺满
    slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)
    print(f"  ✅ 第 {{i+1}} 页已嵌入 PPT")

prs.save(OUTPUT)
print(f"\\n📊 PPT 已生成: {{OUTPUT}} ({{len(screenshots)}} 页)")
'''


def generate_render_script(ppt_data: dict, output_filename: str) -> str:
    """
    生成自包含的 Python 渲染脚本（旧方案，python-pptx 直接渲染）。

    保留用于 fallback 或不需要 playwright 的场景。
    """
    theme_name = ppt_data.get("theme", "tech_blue")
    theme = get_theme(theme_name)

    clean_data = _clean_for_serialization(ppt_data)
    data_json = json.dumps(clean_data, ensure_ascii=False, indent=2)

    colors = {
        "primary": theme.colors.primary,
        "secondary": theme.colors.secondary,
        "accent": theme.colors.accent,
        "bg": theme.colors.bg,
        "text": theme.colors.text,
        "text_light": theme.colors.text_light,
    }
    colors_json = json.dumps(colors)
    title_font = theme.title_font
    body_font = theme.body_font

    return f'''#!/usr/bin/env python3
"""自动生成的 PPT 渲染脚本 — 在沙箱中执行（旧方案 fallback）"""
import json, io, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE

# ── 数据 ──
PPT_DATA = json.loads("""{data_json}""")
COLORS = json.loads('{colors_json}')
TITLE_FONT = "{title_font}"
BODY_FONT = "{body_font}"
OUTPUT = "{output_filename}"

def hex2rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_textbox(slide, left, top, w, h, text, size=18, color="", bold=False, align=PP_ALIGN.LEFT, font=""):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font or BODY_FONT
    if color:
        r.font.color.rgb = hex2rgb(color)
    return tb

def add_bullets(slide, left, top, w, h, items, size=16, color=""):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    c = color or COLORS["text"]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = "  " + str(item)
        r.font.size = Pt(size)
        r.font.name = BODY_FONT
        r.font.color.rgb = hex2rgb(c)

def accent_line(slide, left, top, w):
    s = slide.shapes.add_shape(1, left, top, w, Pt(3))
    s.fill.solid()
    s.fill.fore_color.rgb = hex2rgb(COLORS["accent"])
    s.line.fill.background()

def set_bg(slide, color=""):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex2rgb(color or COLORS["bg"])

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

# ── 布局渲染 ──

def render_title(prs, d):
    s = blank(prs); set_bg(s)
    accent_line(s, Inches(1.5), Inches(3.2), Inches(3))
    add_textbox(s, Inches(1.5), Inches(1.8), Inches(7), Inches(1.5),
                d.get("title",""), size=36, bold=True, color=COLORS["primary"], font=TITLE_FONT)
    if d.get("subtitle"):
        add_textbox(s, Inches(1.5), Inches(3.5), Inches(7), Inches(0.8),
                    d["subtitle"], size=18, color=COLORS["secondary"])

def render_content(prs, d):
    s = blank(prs); set_bg(s)
    add_textbox(s, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.8),
                d.get("title",""), size=28, bold=True, color=COLORS["primary"], font=TITLE_FONT)
    accent_line(s, Inches(0.8), Inches(1.15), Inches(1.5))
    if d.get("bullets"):
        add_bullets(s, Inches(0.8), Inches(1.5), Inches(8.4), Inches(5), d["bullets"], size=18)

def render_two_column(prs, d):
    s = blank(prs); set_bg(s)
    add_textbox(s, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.8),
                d.get("title",""), size=28, bold=True, color=COLORS["primary"], font=TITLE_FONT)
    accent_line(s, Inches(0.8), Inches(1.15), Inches(1.5))
    L, R = d.get("left",{{}}), d.get("right",{{}})
    if L.get("heading"):
        add_textbox(s, Inches(0.8), Inches(1.5), Inches(4), Inches(0.5),
                    L["heading"], size=20, bold=True, color=COLORS["accent"])
    if L.get("bullets"):
        add_bullets(s, Inches(0.8), Inches(2.1), Inches(4), Inches(4.5), L["bullets"], size=16)
    if R.get("heading"):
        add_textbox(s, Inches(5.2), Inches(1.5), Inches(4), Inches(0.5),
                    R["heading"], size=20, bold=True, color=COLORS["accent"])
    if R.get("bullets"):
        add_bullets(s, Inches(5.2), Inches(2.1), Inches(4), Inches(4.5), R["bullets"], size=16)

def render_image_text(prs, d):
    import os
    s = blank(prs); set_bg(s)
    add_textbox(s, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.8),
                d.get("title",""), size=28, bold=True, color=COLORS["primary"], font=TITLE_FONT)
    img_file = d.get("_sandbox_image", "")
    if img_file and os.path.exists(img_file):
        try:
            s.shapes.add_picture(img_file, Inches(0.8), Inches(1.4), Inches(4.2), Inches(4.8))
            print(f"    📷 已嵌入配图: {{img_file}}")
        except Exception as e:
            print(f"    ⚠️ 配图嵌入失败: {{e}}")
            img_file = ""
    if not img_file or not os.path.exists(img_file if img_file else "/nonexist"):
        sh = s.shapes.add_shape(1, Inches(0.8), Inches(1.4), Inches(4.2), Inches(4.8))
        sh.fill.solid()
        sh.fill.fore_color.rgb = hex2rgb("#F0F0F0")
        sh.line.color.rgb = hex2rgb("#E0E0E0")
        tf = sh.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = "[" + d.get("image_prompt","图片") + "]"
        r.font.size = Pt(12); r.font.color.rgb = hex2rgb("#999999")
    if d.get("text"):
        add_textbox(s, Inches(5.3), Inches(1.4), Inches(4), Inches(4.8),
                    d["text"], size=16, color=COLORS["text"])

def render_chart(prs, d):
    s = blank(prs); set_bg(s)
    add_textbox(s, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.8),
                d.get("title",""), size=28, bold=True, color=COLORS["primary"], font=TITLE_FONT)
    cd_raw = d.get("chart_data", {{}})
    cats = cd_raw.get("categories", [])
    series = cd_raw.get("series", [])
    if not cats or not series:
        add_textbox(s, Inches(2), Inches(3), Inches(6), Inches(1),
                    "[图表数据不足]", size=16, color="#999", align=PP_ALIGN.CENTER)
        return
    ct_map = {{"bar": XL_CHART_TYPE.COLUMN_CLUSTERED, "pie": XL_CHART_TYPE.PIE, "line": XL_CHART_TYPE.LINE_MARKERS}}
    ct = ct_map.get(cd_raw.get("type","bar"), XL_CHART_TYPE.COLUMN_CLUSTERED)
    from pptx.chart.data import CategoryChartData
    cd = CategoryChartData(); cd.categories = cats
    for sr in series:
        cd.add_series(sr.get("name",""), sr.get("values",[]))
    s.shapes.add_chart(ct, Inches(1), Inches(1.5), Inches(8), Inches(5), cd)

def render_quote(prs, d):
    s = blank(prs); set_bg(s)
    q = d.get("quote", d.get("text",""))
    add_textbox(s, Inches(1.5), Inches(1.5), Inches(1), Inches(1),
                "\\u201C", size=72, color=COLORS["accent"], bold=True)
    add_textbox(s, Inches(1.5), Inches(2.5), Inches(7), Inches(2.5),
                q, size=24, color=COLORS["text"])
    if d.get("author"):
        add_textbox(s, Inches(1.5), Inches(5.2), Inches(7), Inches(0.5),
                    "—— " + d["author"], size=16, color=COLORS["secondary"], align=PP_ALIGN.RIGHT)

def render_ending(prs, d):
    s = blank(prs); set_bg(s)
    accent_line(s, Inches(3.5), Inches(3), Inches(3))
    add_textbox(s, Inches(1), Inches(2), Inches(8), Inches(1),
                d.get("title","Thanks"), size=44, bold=True, color=COLORS["primary"],
                font=TITLE_FONT, align=PP_ALIGN.CENTER)
    if d.get("subtitle"):
        add_textbox(s, Inches(1), Inches(3.3), Inches(8), Inches(0.6),
                    d["subtitle"], size=16, color=COLORS["secondary"], align=PP_ALIGN.CENTER)

RENDERERS = {{
    "title": render_title, "content": render_content, "two_column": render_two_column,
    "image_text": render_image_text, "chart": render_chart, "quote": render_quote, "ending": render_ending,
}}

# ── 主流程 ──
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

slides = PPT_DATA.get("slides", [])
for i, sd in enumerate(slides):
    layout = sd.get("layout", "content")
    fn = RENDERERS.get(layout, render_content)
    try:
        fn(prs, sd)
        print(f"  ✅ 第 {{i+1}} 页: {{sd.get('title','')}} ({{layout}})")
    except Exception as e:
        print(f"  ⚠️ 第 {{i+1}} 页渲染失败: {{e}}")
        render_content(prs, {{"title": sd.get("title","渲染失败"), "bullets": [str(e)]}})

prs.save(OUTPUT)
print(f"\\n📊 PPT 已生成: {{OUTPUT}} ({{len(slides)}} 页)")
'''


def _clean_for_serialization(data: dict) -> dict:
    """移除不可序列化的字段（如 image_data bytes）。"""
    clean = {}
    for k, v in data.items():
        if isinstance(v, bytes):
            continue
        if isinstance(v, dict):
            clean[k] = _clean_for_serialization(v)
        elif isinstance(v, list):
            clean[k] = [_clean_for_serialization(item) if isinstance(item, dict) else item
                        for item in v if not isinstance(item, bytes)]
        else:
            clean[k] = v
    return clean
